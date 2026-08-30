"""
文件分析工具 — 所有配置值和用户消息从 config 读取
"""

from __future__ import annotations

import asyncio
import logging
import os
from pathlib import Path
from typing import Any, Dict

from config.config_utils import get_file_analysis_config, get_text_message
from webnet.ToolNet.base import BaseTool

logger = logging.getLogger(__name__)

try:
    import pymupdf

    HAS_PDF = True
except ImportError:
    HAS_PDF = False
try:
    from docx import Document

    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False
try:
    from openpyxl import load_workbook

    HAS_XLSX = True
except ImportError:
    HAS_XLSX = False
try:
    from pptx import Presentation

    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
try:
    from PIL import Image

    HAS_PIL = True
except ImportError:
    HAS_PIL = False
try:
    import py7zr

    HAS_7Z = True
except ImportError:
    HAS_7Z = False
try:
    import rarfile

    HAS_RAR = True
except ImportError:
    HAS_RAR = False

FILE_SIGNATURES = {
    b"\x25\x50\x44\x46": ("PDF", ".pdf"),
    b"\x50\x4b\x03\x04": ("ZIP/DOCX/XLSX/PPTX", ".zip"),
    b"\xd0\xcf\x11\xe0": ("DOC/XLS/PPT (OLE)", ".doc"),
    b"\x89\x50\x4e\x47": ("PNG 图片", ".png"),
    b"\xff\xd8\xff": ("JPEG 图片", ".jpg"),
    b"\x47\x49\x46\x38": ("GIF 图片", ".gif"),
    b"\x52\x61\x72\x21": ("RAR 压缩包", ".rar"),
    b"\x37\x7a\xbc\xaf": ("7Z 压缩包", ".7z"),
}

TEXT_EXTENSIONS = {
    ".txt",
    ".md",
    ".py",
    ".js",
    ".ts",
    ".jsx",
    ".tsx",
    ".json",
    ".xml",
    ".yaml",
    ".yml",
    ".toml",
    ".ini",
    ".cfg",
    ".conf",
    ".css",
    ".html",
    ".htm",
    ".csv",
    ".log",
    ".env",
    ".sh",
    ".bat",
    ".ps1",
    ".c",
    ".cpp",
    ".h",
    ".hpp",
    ".java",
    ".go",
    ".rs",
    ".rb",
    ".php",
    ".swift",
    ".kt",
    ".scala",
    ".r",
    ".sql",
    ".lua",
}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp", ".ico", ".tiff"}
CODE_EXTENSIONS = {
    ".py",
    ".js",
    ".ts",
    ".jsx",
    ".tsx",
    ".c",
    ".cpp",
    ".h",
    ".hpp",
    ".java",
    ".go",
    ".rs",
    ".rb",
    ".php",
    ".swift",
    ".kt",
    ".scala",
    ".r",
    ".sql",
    ".lua",
    ".sh",
    ".bat",
    ".ps1",
}


def _limits() -> dict:
    return get_file_analysis_config("limits", default={})


def extract_pdf_pages(file_path: str) -> list:
    """提取 PDF 每页文本 — 弥娅唯一 PDF 文本提取实现（pymupdf 优先，PyPDF2 兜底）

    其他模块（office/pdf_docx_processor、office/invoice_parser、qq/qq_file_reader）
    一律复用此函数，避免重复维护 PDF 解析逻辑。
    """
    if HAS_PDF:
        try:
            doc = pymupdf.open(file_path)
            try:
                return [doc[i].get_text() for i in range(len(doc))]
            finally:
                doc.close()
        except Exception as e:
            logger.warning(f"pymupdf 提取 PDF 失败，回退 PyPDF2: {e}")
    try:
        from PyPDF2 import PdfReader

        reader = PdfReader(file_path)
        return [page.extract_text() or "" for page in reader.pages]
    except Exception as e:
        logger.error(f"PDF 文本提取失败: {e}")
        return []


def extract_pdf_text(file_path: str) -> str:
    """提取 PDF 完整文本 — 唯一实现"""
    return "\n".join(extract_pdf_pages(file_path))


def extract_pdf_metadata(file_path: str) -> Dict[str, str]:
    """提取 PDF 元数据（标题/作者/创建日期/修改日期）"""
    try:
        from PyPDF2 import PdfReader

        metadata = PdfReader(file_path).metadata
        if not metadata:
            return {}
        return {
            "title": str(metadata.get("/Title") or ""),
            "author": str(metadata.get("/Author") or ""),
            "creation_date": str(metadata.get("/CreationDate") or ""),
            "mod_date": str(metadata.get("/ModDate") or ""),
        }
    except Exception:
        return {}


def detect_file_type(file_path: str) -> Dict[str, Any]:
    path = Path(file_path)
    ext = path.suffix.lower()
    if not path.exists():
        return {"type": "unknown", "extension": ext, "description": "文件不存在"}
    size = path.stat().st_size
    try:
        with open(file_path, "rb") as f:
            header = f.read(8)
    except Exception:
        header = b""
    for sig, (name, sig_ext) in FILE_SIGNATURES.items():
        if header.startswith(sig):
            ext = ext or sig_ext
            break
    if ext in TEXT_EXTENSIONS:
        return {
            "type": "text",
            "extension": ext,
            "description": f"文本文件 ({ext})",
            "size": size,
            "is_code": ext in CODE_EXTENSIONS,
        }
    elif ext in IMAGE_EXTENSIONS:
        return {"type": "image", "extension": ext, "description": f"图片文件 ({ext})", "size": size}
    for e2, t2 in [
        (".pdf", "pdf"),
        (".docx", "docx"),
        (".doc", "doc"),
        (".xlsx", "excel"),
        (".xls", "excel"),
        (".pptx", "ppt"),
        (".ppt", "ppt"),
        (".zip", "archive"),
        (".7z", "archive"),
        (".rar", "archive"),
    ]:
        if ext == e2:
            desc = {
                "pdf": "PDF 文档",
                "docx": "Word 文档 (.docx)",
                "doc": "Word 文档 (.doc)",
                "excel": "Excel 表格",
                "ppt": "PowerPoint",
                "archive": "压缩包",
            }.get(t2, t2)
            return {"type": t2, "extension": ext, "description": f"{desc} ({ext})", "size": size}
    return {"type": "unknown", "extension": ext, "description": f"未知类型 ({ext})", "size": size}


class FileAnalyzer:
    @staticmethod
    async def analyze_pdf(file_path: str) -> str:
        if not HAS_PDF:
            return get_text_message("file_analysis", "pdf_install_hint")
        mp = _limits().get("max_pdf_pages", 30)

        def _run():
            pages_text = extract_pdf_pages(file_path)
            total = len(pages_text)
            lines = [get_text_message("file_analysis", "pdf_header", total=total)]
            pages = min(mp, total)
            for i in range(pages):
                text = pages_text[i]
                if text.strip():
                    lines.append(get_text_message("file_analysis", "pdf_page", page=i + 1))
                    lines.append(text[:1500])
            if total > mp:
                lines.append(get_text_message("file_analysis", "pdf_omitted", count=total - mp))
            return "\n".join(lines)

        return await asyncio.to_thread(_run)

    @staticmethod
    async def analyze_docx(file_path: str) -> str:
        if not HAS_DOCX:
            return get_text_message("file_analysis", "docx_install_hint")
        max_chars = _limits().get("max_content_length", 8000)

        def _run():
            doc = Document(file_path)
            lines = [get_text_message("file_analysis", "docx_header", count=len(doc.paragraphs))]
            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    prefix = "## " if para.style.name.startswith("Heading") else ""
                    lines.append(f"\n{prefix}{text}")
                if len("\n".join(lines)) > max_chars:
                    lines.append(get_text_message("file_analysis", "content_truncated"))
                    break
            tables = doc.tables
            if tables:
                lines.append(get_text_message("file_analysis", "docx_tables", count=len(tables)))
            return "\n".join(lines)

        return await asyncio.to_thread(_run)

    @staticmethod
    async def analyze_xlsx(file_path: str) -> str:
        if not HAS_XLSX:
            return get_text_message("file_analysis", "xlsx_install_hint")
        ms = _limits().get("max_xlsx_sheets", 5)
        mr = _limits().get("max_xlsx_rows", 30)

        def _run():
            wb = load_workbook(file_path, read_only=True, data_only=True)
            lines = [
                get_text_message(
                    "file_analysis", "xlsx_header", count=len(wb.sheetnames), sheets=", ".join(wb.sheetnames[:10])
                )
            ]
            for sn in wb.sheetnames[:ms]:
                ws = wb[sn]
                lines.append(get_text_message("file_analysis", "xlsx_sheet", name=sn))
                rows = list(ws.iter_rows(values_only=True))
                for row in rows[:mr]:
                    row_str = " | ".join(str(c) if c is not None else "" for c in row)
                    if row_str.strip():
                        lines.append(row_str)
                if len(rows) > mr:
                    lines.append(get_text_message("file_analysis", "xlsx_omitted", count=len(rows) - mr))
            wb.close()
            return "\n".join(lines)

        return await asyncio.to_thread(_run)

    @staticmethod
    async def analyze_pptx(file_path: str) -> str:
        if not HAS_PPTX:
            return get_text_message("file_analysis", "pptx_install_hint")
        ms = _limits().get("max_ppt_slides", 20)

        def _run():
            prs = Presentation(file_path)
            lines = [get_text_message("file_analysis", "pptx_header", count=len(prs.slides))]
            for i, slide in enumerate(prs.slides):
                if i >= ms:
                    lines.append(get_text_message("file_analysis", "pptx_omitted", count=len(prs.slides) - ms))
                    break
                lines.append(get_text_message("file_analysis", "pptx_slide", num=i + 1))
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                lines.append(text)
            return "\n".join(lines)

        return await asyncio.to_thread(_run)

    @staticmethod
    async def analyze_text(file_path: str) -> str:
        mt = _limits().get("max_text_lines", 200)

        def _run():
            try:
                text = Path(file_path).read_text(encoding="utf-8")
            except UnicodeDecodeError:
                try:
                    text = Path(file_path).read_text(encoding="gbk")
                except Exception:
                    return get_text_message("file_analysis", "text_unreadable")
            all_lines = text.split("\n")
            total = len(all_lines)
            ext = Path(file_path).suffix.lower()
            result = [get_text_message("file_analysis", "text_header", name=Path(file_path).name, lines=total)]
            if ext in CODE_EXTENSIONS:
                result.append(get_text_message("file_analysis", "text_code_type", ext=ext))
            result.append("\n".join(all_lines[:mt]))
            if total > mt:
                result.append(get_text_message("file_analysis", "text_omitted", count=total - mt))
            return "\n".join(result)

        return await asyncio.to_thread(_run)

    @staticmethod
    async def analyze_image(file_path: str) -> str:
        if not HAS_PIL:
            return get_text_message("file_analysis", "image_install_hint")

        def _run():
            img = Image.open(file_path)
            name = Path(file_path).name
            lines = [get_text_message("file_analysis", "image_info", name=name)]
            lines.append(get_text_message("file_analysis", "image_format", format=img.format or "未知"))
            lines.append(get_text_message("file_analysis", "image_size", width=img.size[0], height=img.size[1]))
            lines.append(get_text_message("file_analysis", "image_mode", mode=img.mode))
            if hasattr(img, "info") and img.info and img.info.get("dpi"):
                lines.append(get_text_message("file_analysis", "image_dpi", dpi=img.info["dpi"]))
            fs = os.path.getsize(file_path)
            if fs > 1024 * 1024:
                lines.append(get_text_message("file_analysis", "image_file_size_mb", size=fs / (1024 * 1024)))
            else:
                lines.append(get_text_message("file_analysis", "image_file_size_kb", size=fs / 1024))
            return "\n".join(lines)

        return await asyncio.to_thread(_run)

    @staticmethod
    async def analyze_archive(file_path: str) -> str:
        ext = Path(file_path).suffix.lower()
        max_items = _limits().get("max_archive_files", 50)
        items = []
        total = 0
        if ext == ".zip":
            import zipfile

            with zipfile.ZipFile(file_path, "r") as zf:
                total = len(zf.namelist())
                items = zf.namelist()[:max_items]
        elif ext == ".7z" and HAS_7Z:
            with py7zr.SevenZipFile(file_path, "r") as zf:
                total = len(zf.getnames())
                items = zf.getnames()[:max_items]
        elif ext == ".rar" and HAS_RAR:
            with rarfile.RarFile(file_path, "r") as rf:
                total = len(rf.namelist())
                items = rf.namelist()[:max_items]
        else:
            return get_text_message("file_analysis", "archive_unsupported", ext=ext)
        lines = [get_text_message("file_analysis", "archive_header", name=Path(file_path).name, total=total)]
        for item in items:
            icon = "📁" if item.endswith("/") else "📄"
            lines.append(f"  {icon} {item}")
        if total > max_items:
            lines.append(get_text_message("file_analysis", "archive_omitted", count=total - max_items))
        return "\n".join(lines)

    @staticmethod
    async def analyze(file_path: str) -> str:
        info = detect_file_type(file_path)
        ftype = info["type"]
        try:
            handlers = {
                "pdf": FileAnalyzer.analyze_pdf,
                "docx": FileAnalyzer.analyze_docx,
                "excel": FileAnalyzer.analyze_xlsx,
                "ppt": FileAnalyzer.analyze_pptx,
                "text": FileAnalyzer.analyze_text,
                "image": FileAnalyzer.analyze_image,
                "archive": FileAnalyzer.analyze_archive,
            }
            handler = handlers.get(ftype)
            if handler:
                return await handler(file_path)
            return get_text_message("file_analysis", "type_unsupported", type=info["description"])
        except Exception as e:
            logger.error(f"文件分析失败: {e}", exc_info=True)
            return get_text_message("file_analysis", "analyze_failed", error=str(e)[:200])


class FileAnalysisTool(BaseTool):
    @property
    def config(self) -> Dict[str, Any]:
        return {
            "name": "analyze_file",
            "description": "分析文件内容。支持PDF、Word(.docx)、Excel(.xlsx)、PPT(.pptx)、文本文件、代码文件、图片和压缩包。当用户发送文件、要求分析文件内容、或提到'看看这个文件'时使用。",
            "parameters": {
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "文件路径（本地绝对路径）"},
                    "file_type": {
                        "type": "string",
                        "description": "可选：明确指定文件类型",
                        "enum": ["pdf", "docx", "excel", "ppt", "text", "image", "archive", "auto"],
                        "default": "auto",
                    },
                },
                "required": ["file_path"],
            },
        }

    async def execute(self, args: Dict[str, Any], context) -> str:
        file_path = args.get("file_path", "")
        file_type = args.get("file_type", "auto")
        if not file_path:
            return get_text_message("file_analysis", "path_required")
        if not os.path.exists(file_path):
            return get_text_message("file_analysis", "file_not_found", path=file_path)
        if file_type != "auto":
            handlers = {
                "pdf": FileAnalyzer.analyze_pdf,
                "docx": FileAnalyzer.analyze_docx,
                "excel": FileAnalyzer.analyze_xlsx,
                "ppt": FileAnalyzer.analyze_pptx,
                "text": FileAnalyzer.analyze_text,
                "image": FileAnalyzer.analyze_image,
                "archive": FileAnalyzer.analyze_archive,
            }
            h = handlers.get(file_type)
            return await h(file_path) if h else f"不支持的类型: {file_type}"
        return await FileAnalyzer.analyze(file_path)


class DetectFileTypeTool(BaseTool):
    @property
    def config(self) -> Dict[str, Any]:
        return {
            "name": "detect_file_type",
            "description": "检测文件类型，返回文件的类型、扩展名、大小等信息。当你需要确认文件格式时使用。",
            "parameters": {
                "type": "object",
                "properties": {"file_path": {"type": "string", "description": "文件路径（本地绝对路径）"}},
                "required": ["file_path"],
            },
        }

    async def execute(self, args: Dict[str, Any], context) -> str:
        file_path = args.get("file_path", "")
        if not file_path:
            return get_text_message("file_analysis", "path_required")
        info = detect_file_type(file_path)
        size_str = ""
        if info.get("size"):
            s = info["size"]
            size_str = f" ({s / (1024 * 1024):.1f} MB)" if s > 1024 * 1024 else f" ({s / 1024:.1f} KB)"
        code_note = " [代码文件]" if info.get("is_code") else ""
        return f"类型: {info['description']}{code_note}{size_str}"
