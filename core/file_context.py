"""
统一文件上下文 (FileContext) 数据类

弥娅跨平台文件处理的标准数据模型。
所有平台适配器接收到的文件/图片/语音/视频消息，
统一转换为 FileContext 后传入 DecisionHub。

文件类型枚举:
    image, document, audio, video, archive, code, text, unknown
"""

from __future__ import annotations

import logging
import os
from dataclasses import dataclass, field
from enum import Enum
from pathlib import Path
from typing import Any, Dict, List, Optional

logger = logging.getLogger("Miya.FileContext")


def get_downloads_dir() -> str:
    """获取统一文件下载目录: data/downloads/"""
    project_root = Path(__file__).resolve().parent.parent
    downloads = project_root / "data" / "downloads"
    os.makedirs(downloads, exist_ok=True)
    return str(downloads)


class FileType(str, Enum):
    IMAGE = "image"
    DOCUMENT = "document"
    AUDIO = "audio"
    VIDEO = "video"
    ARCHIVE = "archive"
    CODE = "code"
    TEXT = "text"
    UNKNOWN = "unknown"


class DownloadStatus(str, Enum):
    PENDING = "pending"
    DOWNLOADING = "downloading"
    DONE = "done"
    FAILED = "failed"


@dataclass
class FileContext:
    """统一文件上下文"""

    file_type: str = FileType.UNKNOWN
    file_name: str = ""
    file_url: Optional[str] = None
    file_path: Optional[str] = None
    file_size: int = 0
    mime_type: str = ""
    file_id: str = ""
    file_data: Optional[bytes] = field(default=None, repr=False)
    metadata: Dict[str, Any] = field(default_factory=dict)
    download_status: str = DownloadStatus.PENDING
    analysis_result: Optional[str] = None
    thumbnail_url: Optional[str] = None

    @property
    def is_downloaded(self) -> bool:
        return self.file_path is not None and self.download_status == DownloadStatus.DONE

    @property
    def extension(self) -> str:
        if self.file_name and "." in self.file_name:
            return self.file_name.rsplit(".", 1)[-1].lower()
        return ""

    @classmethod
    def from_image(
        cls,
        url: str = "",
        file_name: str = "",
        file_path: str = "",
        file_id: str = "",
        file_data: Optional[bytes] = None,
        file_size: int = 0,
        mime_type: str = "",
        **kwargs,
    ) -> FileContext:
        return cls(
            file_type=FileType.IMAGE,
            file_url=url,
            file_name=file_name,
            file_path=file_path,
            file_id=file_id,
            file_data=file_data,
            file_size=file_size,
            mime_type=mime_type or "image/unknown",
            metadata=kwargs,
        )

    @classmethod
    def from_file(
        cls,
        url: str = "",
        file_name: str = "",
        file_path: str = "",
        file_id: str = "",
        file_data: Optional[bytes] = None,
        file_size: int = 0,
        mime_type: str = "",
        **kwargs,
    ) -> FileContext:
        file_type = cls._detect_type(file_name, mime_type)
        return cls(
            file_type=file_type,
            file_url=url,
            file_name=file_name,
            file_path=file_path,
            file_id=file_id,
            file_data=file_data,
            file_size=file_size,
            mime_type=mime_type,
            metadata=kwargs,
        )

    @classmethod
    def from_voice(
        cls,
        url: str = "",
        file_name: str = "",
        file_path: str = "",
        file_id: str = "",
        file_data: Optional[bytes] = None,
        file_size: int = 0,
        **kwargs,
    ) -> FileContext:
        return cls(
            file_type=FileType.AUDIO,
            file_url=url,
            file_name=file_name,
            file_path=file_path,
            file_id=file_id,
            file_data=file_data,
            file_size=file_size,
            mime_type="audio/unknown",
            metadata=kwargs,
        )

    @classmethod
    def from_video(
        cls,
        url: str = "",
        file_name: str = "",
        file_path: str = "",
        file_id: str = "",
        file_size: int = 0,
        **kwargs,
    ) -> FileContext:
        return cls(
            file_type=FileType.VIDEO,
            file_url=url,
            file_name=file_name,
            file_path=file_path,
            file_id=file_id,
            file_size=file_size,
            mime_type="video/unknown",
            metadata=kwargs,
        )

    @staticmethod
    def _detect_type(file_name: str, mime_type: str = "") -> str:
        if mime_type:
            if mime_type.startswith("image/"):
                return FileType.IMAGE
            if mime_type.startswith("audio/"):
                return FileType.AUDIO
            if mime_type.startswith("video/"):
                return FileType.VIDEO

        ext = file_name.rsplit(".", 1)[-1].lower() if "." in file_name else ""
        image_exts = {"jpg", "jpeg", "png", "gif", "webp", "bmp", "svg", "ico", "tiff"}
        audio_exts = {"mp3", "wav", "ogg", "aac", "flac", "m4a", "wma", "opus", "amr"}
        video_exts = {"mp4", "avi", "mkv", "mov", "wmv", "flv", "webm", "3gp"}
        archive_exts = {"zip", "rar", "7z", "tar", "gz", "bz2", "xz"}
        code_exts = {"py", "js", "ts", "java", "cpp", "c", "go", "rs", "php", "rb", "sh", "bat", "ps1", "sql"}
        text_exts = {"txt", "log", "md", "json", "xml", "html", "csv", "yml", "yaml", "ini", "cfg", "conf", "pdf"}

        if ext in image_exts:
            return FileType.IMAGE
        if ext in audio_exts:
            return FileType.AUDIO
        if ext in video_exts:
            return FileType.VIDEO
        if ext in archive_exts:
            return FileType.ARCHIVE
        if ext in code_exts:
            return FileType.CODE
        if ext in text_exts:
            return FileType.TEXT
        return FileType.DOCUMENT

    def to_dict(self) -> Dict[str, Any]:
        return {
            "file_type": self.file_type.value if isinstance(self.file_type, FileType) else self.file_type,
            "file_name": self.file_name,
            "file_url": self.file_url,
            "file_path": self.file_path,
            "file_size": self.file_size,
            "mime_type": self.mime_type,
            "file_id": self.file_id,
            "download_status": self.download_status.value
            if isinstance(self.download_status, DownloadStatus)
            else self.download_status,
            "analysis_result": self.analysis_result,
            "metadata": self.metadata,
        }

    def analyze_content(self) -> Optional[str]:
        """解析文件内容提取文本（支持 docx/pdf/txt/code 等常见格式）"""
        if self.analysis_result:
            return self.analysis_result

        if not self.file_data:
            if self.file_path and os.path.exists(self.file_path):
                with open(self.file_path, "rb") as f:
                    self.file_data = f.read()
                self.file_size = len(self.file_data)
            else:
                return None

        result = None
        ext = self.extension

        if ext == "docx":
            result = self._read_docx()
        elif ext == "doc":
            result = self._read_doc()
        elif ext == "pdf":
            result = self._read_pdf()
        elif ext in ("xlsx", "xls"):
            result = self._read_xlsx()
        elif ext in ("pptx", "ppt"):
            result = self._read_pptx()
        elif ext in (
            "txt",
            "md",
            "log",
            "json",
            "xml",
            "yaml",
            "yml",
            "toml",
            "ini",
            "cfg",
            "conf",
            "csv",
            "html",
            "css",
            "py",
            "js",
            "ts",
            "java",
            "go",
            "rs",
            "c",
            "cpp",
            "h",
            "sh",
            "bat",
            "ps1",
            "rb",
            "php",
            "lua",
            "r",
            "sql",
        ):
            result = self._read_text()

        if result:
            self.analysis_result = result
        return result

    def _read_docx(self) -> Optional[str]:
        try:
            from io import BytesIO
            from docx import Document

            doc = Document(BytesIO(self.file_data))
            lines = []
            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    prefix = "# " if para.style and para.style.name and para.style.name.startswith("Heading") else ""
                    lines.append(f"{prefix}{text}")
            result = "\n".join(lines)
            return result[:8000] if len(result) > 8000 else result
        except Exception as e:
            logger.debug(f"docx 解析失败: {e}")
            return None

    def _read_doc(self) -> Optional[str]:
        """从旧 .doc (OLE2 二进制) 格式中提取文本

        策略：OLE2 容器中 Word 文档文本可能以 UTF-16LE 存储，
        优先搜索可读的 Unicode 文本块。
        """
        import re

        try:
            data = self.file_data
            blocks = []

            # 策略1: 尝试在二进制中查找 UTF-16LE 文本段
            try:
                decoded = data.decode("utf-16-le", errors="ignore")
                cleaned = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f]", "", decoded)
                cleaned = re.sub(
                    r"[^\u4e00-\u9fff\u3000-\u303f\uff00-\uffef"
                    r"a-zA-Z0-9\s.,;:!?()\[\]{}\'\"\-_+=@#$%^&*|\\/<>\n\r\t\u2000-\u206f]",
                    " ",
                    cleaned,
                )
                cleaned = re.sub(r"\s{2,}", "\n", cleaned).strip()
                if len(cleaned) > 100:
                    blocks.append(cleaned)
            except Exception:
                pass

            # 策略2: Latin-1 解码 + 提取可读单词序列
            if not blocks:
                text = data.decode("latin-1", errors="ignore")
                text = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\x9f]", "", text)
                segments = re.findall(r"[\u4e00-\u9fff\u3000-\u303f\uff00-\uffefa-zA-Z0-9\u2000-\u206f]{3,}", text)
                if segments:
                    blocks.append(" ".join(segments))

            if not blocks:
                return None

            result = "\n".join(blocks)
            # 垃圾检测：如果结果中可读文本占比低于 30%，放弃
            readable = sum(1 for c in result if c.isprintable() or c in "\n\r\t")
            if len(result) > 0 and readable / len(result) < 0.3:
                return None
            return result[:8000] if len(result) > 8000 else result
        except Exception as e:
            logger.debug(f"doc 解析失败: {e}")
            return None

    def _read_pdf(self) -> Optional[str]:
        try:
            import tempfile

            tmp = tempfile.NamedTemporaryFile(suffix=".pdf", delete=False)
            try:
                tmp.write(self.file_data)
                tmp.close()
                import pymupdf

                doc = pymupdf.open(tmp.name)
                lines = []
                max_pages = min(len(doc), 30)
                for i in range(max_pages):
                    text = doc[i].get_text()
                    if text.strip():
                        lines.append(f"--- 第{i + 1}页 ---\n{text[:1500]}")
                doc.close()
                return "\n".join(lines) if lines else None
            finally:
                import os as _os

                _os.unlink(tmp.name)
        except Exception as e:
            logger.debug(f"pdf 解析失败: {e}")
            return None

    def _read_xlsx(self) -> Optional[str]:
        try:
            from io import BytesIO
            from openpyxl import load_workbook

            wb = load_workbook(BytesIO(self.file_data), read_only=True, data_only=True)
            lines = [f"Excel 表格, {len(wb.sheetnames)} 个工作表: {', '.join(wb.sheetnames[:10])}"]
            max_sheets = 5
            for sn in wb.sheetnames[:max_sheets]:
                ws = wb[sn]
                rows = list(ws.iter_rows(values_only=True))
                lines.append(f"\n=== {sn} ({len(rows)} 行) ===")
                for row in rows[:30]:
                    row_str = " | ".join(str(c) if c is not None else "" for c in row)
                    if row_str.strip():
                        lines.append(row_str)
            wb.close()
            result = "\n".join(lines)
            return result[:8000] if len(result) > 8000 else result
        except Exception as e:
            logger.debug(f"xlsx 解析失败: {e}")
            return None

    def _read_pptx(self) -> Optional[str]:
        try:
            from io import BytesIO
            from pptx import Presentation

            prs = Presentation(BytesIO(self.file_data))
            lines = [f"PPT 演示, {len(prs.slides)} 页"]
            max_slides = 20
            for i, slide in enumerate(prs.slides):
                if i >= max_slides:
                    break
                lines.append(f"\n--- 第{i + 1}页 ---")
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                lines.append(text)
            result = "\n".join(lines)
            return result[:8000] if len(result) > 8000 else result
        except Exception as e:
            logger.debug(f"pptx 解析失败: {e}")
            return None

    def _read_text(self) -> Optional[str]:
        try:
            text = self.file_data.decode("utf-8")
        except UnicodeDecodeError:
            try:
                text = self.file_data.decode("gbk")
            except Exception:
                return None
        max_len = 8000
        return text[:max_len] if len(text) > max_len else text

    def to_text_summary(self) -> str:
        """生成文件描述文本，注入到 AI 消息上下文中"""
        parts = []
        type_names = {
            FileType.IMAGE: "图片",
            FileType.DOCUMENT: "文档",
            FileType.AUDIO: "语音",
            FileType.VIDEO: "视频",
            FileType.ARCHIVE: "归档文件",
            FileType.CODE: "代码文件",
            FileType.TEXT: "文本文件",
            FileType.UNKNOWN: "文件",
        }
        type_label = type_names.get(self.file_type, "文件")
        parts.append(f"[{type_label}]")
        if self.file_name:
            parts.append(f" 文件名: {self.file_name}")
        if self.file_size > 0:
            if self.file_size < 1024:
                parts.append(f" 大小: {self.file_size}B")
            elif self.file_size < 1024 * 1024:
                parts.append(f" 大小: {self.file_size / 1024:.1f}KB")
            else:
                parts.append(f" 大小: {self.file_size / (1024 * 1024):.1f}MB")

        if self.analysis_result:
            parts.append(f" 分析结果: {self.analysis_result}")

        return "\n".join(parts)


@dataclass
class OutboundFile:
    """出站文件——弥娅要发送给用户的文件

    与 FileContext（入站）对称，OutboundFile 描述弥娅需要发送出去的文件。
    支持三种来源：本地路径、内存字节、远程 URL。
    """

    file_path: str = ""
    file_name: str = ""
    mime_type: str = ""
    caption: str = ""
    file_size: int = 0
    file_data: Optional[bytes] = field(default=None, repr=False)
    source: str = ""  # local_path / bytes / url
    metadata: Dict[str, Any] = field(default_factory=dict)

    @classmethod
    def from_local(cls, path: str, file_name: str = "", caption: str = "", **kwargs) -> OutboundFile:
        if not path:
            raise ValueError("file_path is required for local source")
        fpath = Path(path)
        if not fpath.exists():
            raise FileNotFoundError(f"文件不存在: {path}")
        name = file_name or fpath.name
        size = fpath.stat().st_size
        mime = cls._guess_mime(name)
        file_type = OutboundFile._detect_type(name, mime)
        return cls(
            file_path=str(fpath.resolve()),
            file_name=name,
            mime_type=mime,
            caption=caption,
            file_size=size,
            source="local_path",
            metadata={"file_type": file_type, **kwargs},
        )

    @classmethod
    def from_bytes(cls, data: bytes, filename: str, mime_type: str = "", caption: str = "") -> OutboundFile:
        if not data:
            raise ValueError("file_data cannot be empty")
        name = filename or "file.bin"
        mime = mime_type or cls._guess_mime(name)
        file_type = OutboundFile._detect_type(name, mime)
        return cls(
            file_name=name,
            mime_type=mime,
            caption=caption,
            file_size=len(data),
            file_data=data,
            source="bytes",
            metadata={"file_type": file_type},
        )

    @classmethod
    def from_url(cls, url: str, filename: str = "", caption: str = "") -> OutboundFile:
        if not url:
            raise ValueError("url is required for URL source")
        name = filename or url.rsplit("/", 1)[-1].split("?", 1)[0] or "download.bin"
        mime = cls._guess_mime(name)
        file_type = OutboundFile._detect_type(name, mime)
        return cls(
            file_name=name,
            mime_type=mime,
            caption=caption,
            source="url",
            metadata={"url": url, "file_type": file_type},
        )

    @property
    def is_local(self) -> bool:
        return self.source == "local_path" and bool(self.file_path)

    @property
    def is_bytes(self) -> bool:
        return self.source == "bytes" and bool(self.file_data)

    @property
    def is_url(self) -> bool:
        return self.source == "url"

    @property
    def is_image(self) -> bool:
        if self.mime_type and self.mime_type.startswith("image/"):
            return True
        ft = self.metadata.get("file_type", "")
        return ft == FileType.IMAGE or ft == "image"

    @property
    def extension(self) -> str:
        if self.file_name and "." in self.file_name:
            return self.file_name.rsplit(".", 1)[-1].lower()
        return ""

    @staticmethod
    def _guess_mime(filename: str) -> str:
        ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
        mime_map = {
            "jpg": "image/jpeg",
            "jpeg": "image/jpeg",
            "png": "image/png",
            "gif": "image/gif",
            "webp": "image/webp",
            "bmp": "image/bmp",
            "svg": "image/svg+xml",
            "mp3": "audio/mpeg",
            "wav": "audio/wav",
            "ogg": "audio/ogg",
            "aac": "audio/aac",
            "flac": "audio/flac",
            "opus": "audio/opus",
            "mp4": "video/mp4",
            "avi": "video/avi",
            "mkv": "video/x-matroska",
            "mov": "video/quicktime",
            "webm": "video/webm",
            "pdf": "application/pdf",
            "doc": "application/msword",
            "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "xls": "application/vnd.ms-excel",
            "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "ppt": "application/vnd.ms-powerpoint",
            "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "zip": "application/zip",
            "rar": "application/x-rar-compressed",
            "7z": "application/x-7z-compressed",
            "tar": "application/x-tar",
            "gz": "application/gzip",
            "txt": "text/plain",
            "md": "text/markdown",
            "json": "application/json",
            "xml": "application/xml",
            "html": "text/html",
            "csv": "text/csv",
        }
        return mime_map.get(ext, "application/octet-stream")

    @staticmethod
    def _detect_type(file_name: str, mime_type: str = "") -> str:
        if mime_type:
            if mime_type.startswith("image/"):
                return FileType.IMAGE
            if mime_type.startswith("audio/"):
                return FileType.AUDIO
            if mime_type.startswith("video/"):
                return FileType.VIDEO
        return FileContext._detect_type(file_name, mime_type)

    def validate(self) -> bool:
        if self.source == "local_path":
            return bool(self.file_path) and Path(self.file_path).exists()
        if self.source == "bytes":
            return bool(self.file_data)
        if self.source == "url":
            return bool(self.metadata.get("url"))
        return False

    def to_dict(self) -> Dict[str, Any]:
        return {
            "file_path": self.file_path,
            "file_name": self.file_name,
            "mime_type": self.mime_type,
            "caption": self.caption,
            "file_size": self.file_size,
            "source": self.source,
            "is_image": self.is_image,
            "extension": self.extension,
            "metadata": self.metadata,
        }


def create_file_contexts_from_platform(
    platform_id: str,
    files_raw: list,
) -> List[FileContext]:
    """从平台原始文件数据创建 FileContext 列表 (工厂函数)"""
    results: List[FileContext] = []
    for f in files_raw:
        if isinstance(f, FileContext):
            results.append(f)
        elif isinstance(f, dict):
            results.append(FileContext(**{k: v for k, v in f.items() if k in FileContext.__dataclass_fields__}))
        else:
            logger.warning(f"[{platform_id}] 无法识别的文件数据格式: {type(f)}")
    return results
